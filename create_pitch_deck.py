from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0xE8, 0x4D, 0x0F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF0, 0xF2, 0xF5)
GRAY = RGBColor(0x6C, 0x75, 0x7D)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    if width is None: width = prs.slide_width
    if height is None: height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_bar(slide, left=0, top=0, width=None, height=Inches(0.06)):
    if width is None: width = prs.slide_width
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
        p.level = 0
    return txBox

# === SLIDE 1: TITLE ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, top=Inches(2.8))
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             'Event-Driven Congestion (Planned & Unplanned)', font_size=44, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.8),
             'Real-Time Traffic Incident Response System for Bengaluru Traffic Police', font_size=22, color=RGBColor(0xA0,0xAE,0xC0))
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
             'Powered by BTP ASTraM Historical Data', font_size=16, color=GRAY)
add_text_box(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.7),
             'ASTraM Event Congestion Command', font_size=24, bold=True, color=ACCENT)
add_text_box(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.6),
             'Team: The ML Mavericks', font_size=18, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.4),
             'Predict · Deploy · Learn', font_size=14, color=GRAY)

# === SLIDE 2: PROBLEM ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_shape_bg(slide, DARK, left=0, top=0, width=Inches(5.5), height=prs.slide_height)
add_text_box(slide, Inches(0.6), Inches(0.8), Inches(4.5), Inches(0.8),
             'THE PROBLEM', font_size=14, bold=True, color=ACCENT)
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(4.5), Inches(1.5),
             'Traffic congestion in Bengaluru costs  ₹38,000 Cr annually', font_size=30, bold=True, color=WHITE)
add_bullet_text(slide, Inches(0.6), Inches(3.4), Inches(4.5), Inches(3.5), [
    '▸  22 major corridors, 800+ junctions — manual incident management is slow',
    '▸  No data-driven resource allocation — officers, barricades, patrols deployed ad-hoc',
    '▸  Diversion routes decided on the fly, no alternate route computation',
    '▸  No feedback loop — past responses are not used to improve future decisions',
    '▸  Disconnected systems: incidents logged in ASTraM, response planned manually',
], font_size=16, color=RGBColor(0xC0,0xCC,0xDD))

add_text_box(slide, Inches(6.5), Inches(0.8), Inches(6), Inches(0.8),
             'OUR SOLUTION', font_size=14, bold=True, color=ACCENT)
add_text_box(slide, Inches(6.5), Inches(1.6), Inches(6), Inches(1.5),
             'A unified command dashboard that predicts and optimizes in real time', font_size=28, bold=True, color=DARK)
add_bullet_text(slide, Inches(6.5), Inches(3.4), Inches(6), Inches(3.5), [
    '→  ML models predict severity, closure, duration, and resource needs',
    '→  Corridor-level routing engine computes actual alternate routes',
    '→  Real-time monitoring via simulation or live TomTom traffic API',
    '→  Feedback-driven retraining pipeline improves predictions over time',
    '→  Single dashboard: analysis, monitor, feedback, all in one place',
], font_size=16, color=DARK)

# === SLIDE 3: HOW IT WORKS ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             'SYSTEM ARCHITECTURE', font_size=32, bold=True, color=WHITE)

steps = [
    ('1', 'Event Ingestion', 'Historical ASTraM data (8,200+ events) or real-time feed\n(simulation / TomTom Traffic API)'),
    ('2', 'Feature Engineering', 'Extract temporal, spatial, and categorical features.\nCompute severity labels from closure + duration + cause heuristics.'),
    ('3', 'ML Prediction', '6 Random Forest models predict: closure, duration bucket,\nseverity level, officers, barricades, patrol vehicles.'),
    ('4', 'Resource & Routing', 'Resource optimizer assigns deployment points.\nRouting engine computes alternate routes from corridor graph.'),
    ('5', 'Dashboard & Feedback', 'Results shown on map + analysis panel.\nGround truth corrections fed back into retraining pipeline.'),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.7) + i * Inches(1.1)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, Inches(1.6), y - Inches(0.05), Inches(3), Inches(0.5),
                 title, font_size=18, bold=True, color=ACCENT)
    add_text_box(slide, Inches(4.8), y - Inches(0.05), Inches(7.5), Inches(0.7),
                 desc, font_size=13, color=RGBColor(0xC0,0xCC,0xDD))

# === SLIDE 4: ML MODELS ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             'MACHINE LEARNING MODELS', font_size=32, bold=True, color=DARK)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             '6 scikit-learn models trained on 7,532 processed events from ASTraM historical data', font_size=16, color=GRAY)

models = [
    ('closure_model', 'Random Forest Classifier', 'Predicts if road closure is required', '87.4% accuracy'),
    ('duration_model', 'Random Forest Classifier', '3-class: short / medium / extended (trained on 2,722 actual-duration events)', '59.6% accuracy'),
    ('severity_model', 'Random Forest Classifier', '3-class: High / Medium / Low (excludes leakage features)', '81.6% accuracy'),
    ('officers_model', 'Random Forest Regressor', 'Predicts number of officers needed', 'MAE 2.33'),
    ('barricades_model', 'Random Forest Regressor', 'Predicts number of barricades needed', 'MAE 1.74'),
    ('patrols_model', 'Random Forest Regressor', 'Predicts number of patrol vehicles needed', 'MAE 0.68'),
]
for i, (name, algo, desc, perf) in enumerate(models):
    row = i // 3
    col = i % 3
    x = Inches(0.8) + col * Inches(4.2)
    y = Inches(2.0) + row * Inches(2.5)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.1))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    card.line.color.rgb = RGBColor(0xDD, 0xE0, 0xE7)
    card.line.width = Pt(1)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.4),
                 name, font_size=15, bold=True, color=DARK)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.55), Inches(3.2), Inches(0.3),
                 algo, font_size=11, color=ACCENT)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.2), Inches(0.6),
                 desc, font_size=11, color=DARK)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.55), Inches(3.2), Inches(0.3),
                 perf, font_size=12, bold=True, color=RGBColor(0x27, 0xAE, 0x60))

# === SLIDE 5: FEATURES ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             'KEY FEATURES', font_size=32, bold=True, color=WHITE)

features = [
    ('Interactive Map', 'Select any of 22 Bengaluru corridors.\nVisualize incident location, severity zones,\nand computed diversion routes.'),
    ('Prediction Pipeline', 'Real-time severity scoring, closure\nprobability, duration estimate, and\nresource deployment recommendation.'),
    ('Feedback & Retraining', 'Log every prediction. Edit ground truth.\nRetrain models with one click to\nincorporate real-world outcomes.'),
    ('Live Monitoring', 'Simulated real-time events (or live\nTomTom Traffic API). SSE stream\npushes events to dashboard.'),
    ('Corridor Routing', 'Graph-based routing engine computes\nactual alternate routes using the\nBangalore major arterial network.'),
    ('Resource Optimization', 'Data-driven allocation of officers,\nbarricades, and patrol vehicles\nbased on event profile + historical data.'),
]
for i, (title, desc) in enumerate(features):
    row = i // 3
    col = i % 3
    x = Inches(0.8) + col * Inches(4.2)
    y = Inches(1.7) + row * Inches(2.7)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x23, 0x35, 0x58)
    card.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(3.2), Inches(0.4),
                 title, font_size=17, bold=True, color=ACCENT)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(1.3),
                 desc, font_size=13, color=RGBColor(0xB0,0xBC,0xCF))

# === SLIDE 6: TECH STACK ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             'TECHNOLOGY STACK', font_size=32, bold=True, color=DARK)

techs = [
    ('Backend', ['Flask (Python)', 'RESTful API', 'Server-Sent Events (SSE)', '11 API endpoints + worker thread']),
    ('ML / Data', ['scikit-learn (6 models)', 'pandas / numpy', 'joblib model persistence', 'Random Forest (Classifier + Regressor)']),
    ('Frontend', ['Vanilla JS', 'Leaflet.js (map)', 'Bootstrap-like CSS grid', 'EventSource (SSE client)']),
    ('Infrastructure', ['Runs on any Python 3.10+ host', 'No database required', 'CSV-based persistence', 'Auto-launches browser on start']),
]
for i, (title, items) in enumerate(techs):
    x = Inches(0.8) + i * Inches(3.2)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(2.9), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    card.line.color.rgb = RGBColor(0xDD, 0xE0, 0xE7)
    card.line.width = Pt(1)
    add_text_box(slide, x + Inches(0.25), Inches(1.9), Inches(2.4), Inches(0.4),
                 title, font_size=18, bold=True, color=DARK)
    add_bullet_text(slide, x + Inches(0.25), Inches(2.5), Inches(2.4), Inches(2.5),
                    [f'•  {item}' for item in items], font_size=13, color=DARK)

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
             'Single-file server.py (~650 lines) + modular src/ package (~2,000 lines total)', font_size=13, color=GRAY)

# === SLIDE 7: DEMO SCREENSHOTS PLACEHOLDER ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             'DASHBOARD OVERVIEW', font_size=32, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             'Single-page dashboard with four integrated tabs', font_size=18, color=RGBColor(0xA0,0xAE,0xC0))

panels = [
    ('Analysis Tab', 'Select corridor → event type → get\nseverity, resources, diversion routes.\nMap shows incident + buffer zones.'),
    ('Feedback Log', 'Browse 100+ predictions. Sort by\nseverity, filter by corridor.\nEdit ground truth, retrain models.'),
    ('Live Monitor', 'Start/stop real-time event stream.\nChoose simulation or live TomTom.\nAPI key configurable via UI.'),
    ('Retraining', 'Train resource models from feedback.\nRefresh all models with one click.\nPerformance metrics displayed.'),
]
for i, (title, desc) in enumerate(panels):
    x = Inches(0.8) + i * Inches(3.2)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.4), Inches(2.9), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x23, 0x35, 0x58)
    card.line.fill.background()
    add_text_box(slide, x + Inches(0.25), Inches(2.6), Inches(2.4), Inches(0.4),
                 title, font_size=17, bold=True, color=ACCENT)
    add_text_box(slide, x + Inches(0.25), Inches(3.2), Inches(2.4), Inches(1.5),
                 desc, font_size=13, color=RGBColor(0xB0,0xBC,0xCF))

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
             'Map shows 22 Bangalore corridors with real road-aligned coordinates · Live circle markers color-coded by severity',
             font_size=13, color=GRAY)

# === SLIDE 8: DATA ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             'DATA & INSIGHTS', font_size=32, bold=True, color=DARK)

stats = [
    ('8,200+', 'Historical ASTraM events\nspanning 2018–2024'),
    ('22', 'Major Bengaluru corridors\nmapped with real coordinates'),
    ('96', 'Junctions with\nlocation data'),
    ('7,532', 'Processed feature-engineered\nrows (processed_theme3.csv)'),
    ('105', 'Feedback log entries\nfrom model predictions'),
]
for i, (num, desc) in enumerate(stats):
    x = Inches(0.5) + i * Inches(2.5)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(2.2), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    card.line.color.rgb = RGBColor(0xDD, 0xE0, 0xE7)
    card.line.width = Pt(1)
    add_text_box(slide, x + Inches(0.15), Inches(1.85), Inches(1.9), Inches(0.6),
                 num, font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.5), Inches(1.9), Inches(1.0),
                 desc, font_size=12, color=DARK, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
             'Key data insight: Priority is a deterministic rule (corridor = High, Non-corridor = Low), not a predictive task.',
             font_size=14, bold=True, color=DARK)
add_text_box(slide, Inches(0.8), Inches(4.9), Inches(11), Inches(0.5),
             '80% of duration values are imputed (median 8.12 hrs) — duration model trained only on 2,722 actual-duration events.',
             font_size=14, color=GRAY)

# === SLIDE 9: IMPACT ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             'BUSINESS IMPACT', font_size=32, bold=True, color=WHITE)

impacts = [
    ('Faster Response', 'From minutes to seconds:\ninstant severity assessment\nand resource recommendation.'),
    ('Optimized Resources', 'Data-driven allocation instead\nof ad-hoc deployment.\nRight resources, right place.'),
    ('Continuous Learning', 'Every incident improves the\nmodels. Feedback loop ensures\ngrowing accuracy over time.'),
    ('Unified Platform', 'Analysis, monitoring, feedback,\nand retraining in one dashboard.\nNo context switching.'),
]
for i, (title, desc) in enumerate(impacts):
    x = Inches(0.8) + i * Inches(3.2)
    add_text_box(slide, x, Inches(1.7), Inches(2.9), Inches(0.4),
                 title, font_size=20, bold=True, color=ACCENT)
    add_text_box(slide, x, Inches(2.3), Inches(2.9), Inches(1.5),
                 desc, font_size=15, color=RGBColor(0xC0,0xCC,0xDD))

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
             'Potential to reduce congestion-related economic losses through faster, data-driven incident management',
             font_size=16, color=GRAY)

# === SLIDE 10: THANK YOU ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, top=Inches(2.6))
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
             'Thank You', font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
             'ASTraM Event Congestion Command — Event-Driven Congestion (Planned & Unplanned)', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.6),
             'Built with Flask, scikit-learn, Leaflet.js · 6 ML models · 22 Bangalore corridors',
             font_size=16, color=RGBColor(0xA0,0xAE,0xC0), align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.6),
             'Team: The ML Mavericks', font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
team_members = ['Eaknath M S (Team Leader)', 'Ashish D', 'Sabarishh P', 'Prasanth Kanaga Sabai S']
for i, name in enumerate(team_members):
    add_text_box(slide, Inches(2) + (i % 2) * Inches(5), Inches(5.0 + 0.4 * (i // 2)), Inches(5), Inches(0.4),
                 name, font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

prs.save('ML_Mavericks_ASTraM_Pitch_Deck.pptx')
print('Created: ML_Mavericks_ASTraM_Pitch_Deck.pptx')
