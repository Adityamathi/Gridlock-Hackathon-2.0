from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0F, 0x17, 0x2A)
ROYAL_BLUE = RGBColor(0x25, 0x63, 0xEB)
SKY_BLUE = RGBColor(0x38, 0xBD, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
OFF_WHITE_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
SLATE = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_r(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def add_rr(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def tb(slide, l, t, w, h, txt, sz=18, clr=WHITE, b=False, al=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = b; p.font.name = 'Calibri'; p.alignment = al

def mt(slide, l, t, w, h, lines, sz=13, clr=SLATE):
    bx = slide.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    for i, item in enumerate(lines):
        p = bx.text_frame.paragraphs[0] if i == 0 else bx.text_frame.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.name = 'Calibri'; p.space_after = Pt(4)

def top_bar(slide):
    add_r(slide, Inches(0), Inches(0), Inches(13.333), Pt(5), ROYAL_BLUE)

def slide_title(slide, title, sub=None):
    top_bar(slide)
    tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), title, 28, WHITE, True)
    if sub: tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4), sub, 14, SKY_BLUE)

# ======================== SLIDE 1: TITLE ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_r(slide, Inches(0), Inches(2.8), Inches(13.333), Pt(5), ROYAL_BLUE)
tb(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.0), "Event-Driven Congestion (Planned & Unplanned)", 22, WHITE, al=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.7), "Real-Time Traffic Incident Response System for Bengaluru Traffic Police", 17, WHITE, al=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.4), "Powered by BTP ASTraM Historical Data", 14, SKY_BLUE, al=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6), "Grid Unlock Command", 30, WHITE, True, PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.4), "Platform Updates & Enhancements", 16, SLATE, al=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4), "Team: The ML Mavericks", 16, WHITE, al=PP_ALIGN.CENTER)
tb(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.3), "Eaknath M S (Lead)  |  Ashish D  |  Sabarishh P  |  Prasanth Kanaga Sabai S", 12, SLATE, al=PP_ALIGN.CENTER)

# ======================== SLIDE 2: CORRIDOR EXPANSION ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
slide_title(slide, "CORRIDOR EXPANSION", "Scaling from 21 named corridors to 147 with address-based extraction")
tb(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.3), "Challenge", 16, ROYAL_BLUE, True)
tb(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.4), "Only 22 named corridors were mapped, leaving 6,000+ events on unmapped roads unaddressable by the prediction pipeline.", 13, SLATE)
add_rr(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.6), DARK_CARD)
tb(slide, Inches(1.1), Inches(2.8), Inches(5.0), Inches(0.3), "Before: 22 Named Corridors", 15, WHITE, True)
mt(slide, Inches(1.1), Inches(3.2), Inches(5.0), Inches(2.8), ["21 named corridors manually curated", "No non-corridor road recognition", "Events on unmapped roads fell into generic bucket", "Limited geographic coverage", "Area context missing from predictions"], 13, SLATE)
add_rr(slide, Inches(6.8), Inches(2.6), Inches(5.5), Inches(3.6), DARK_CARD)
tb(slide, Inches(7.1), Inches(2.8), Inches(5.0), Inches(0.3), "After: 147 Total Corridors", 15, WHITE, True)
mt(slide, Inches(7.1), Inches(3.2), Inches(5.0), Inches(2.8), ["562 unique road names extracted from address field", "Filtered: roads with 5+ events (125 roads retained)", "Deduplicated similar names into single entries", "Roads shown as Road Name, Area Name format", "1,145 unmapped rows grouped as Non-corridor"], 13, SLATE)
add_r(slide, Inches(0.8), Inches(6.6), Inches(11.5), Pt(2), ROYAL_BLUE)
tb(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.3), "147 corridors: 21 named + 125 address-extracted + 1 fallback. All models retrained. Centroids, graph nodes, and impact radii updated across 6 modules.", 11, SKY_BLUE)

# ======================== SLIDE 3: EXTRACTION METHODOLOGY ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
top_bar(slide)
tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), "ADDRESS-BASED ROAD EXTRACTION", 28, DARK_TEXT, True)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4), "Systematic pipeline to extract and normalize non-corridor road names from raw address data", 14, GRAY_TEXT)
for i, (num, title, desc) in enumerate([("1","Parse","Extract first two parts from raw address field"),("2","Count","Aggregate by road name, filter for 5+"),("3","Normalize","Merge duplicate names into one entry"),("4","Map","Assign each road to most common area"),("5","Integrate","Add to network with centroids and radii")]):
    x = Inches(0.8 + i * 2.4)
    add_rr(slide, x, Inches(2.0), Inches(2.1), Inches(2.8), OFF_WHITE_BG)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), Inches(2.3), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = ROYAL_BLUE; c.line.fill.background()
    cp = c.text_frame.paragraphs[0]; cp.text = num; cp.font.size = Pt(18); cp.font.color.rgb = WHITE; cp.font.bold = True; cp.font.name = 'Calibri'; cp.alignment = PP_ALIGN.CENTER
    tb(slide, x + Inches(0.15), Inches(3.0), Inches(1.8), Inches(0.3), title, 15, DARK_TEXT, True, PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.15), Inches(3.4), Inches(1.8), Inches(1.2), desc, 11, GRAY_TEXT, al=PP_ALIGN.CENTER)
for i, (v, l) in enumerate([("562","Unique road names"),("125","Roads with 5+ events"),("1,445","Rows remapped"),("7,532","Total processed rows")]):
    x = Inches(0.8 + i * 3.1)
    tb(slide, x, Inches(5.3), Inches(2.8), Inches(0.4), v, 26, ROYAL_BLUE, True, PP_ALIGN.CENTER)
    tb(slide, x, Inches(5.7), Inches(2.8), Inches(0.3), l, 12, GRAY_TEXT, al=PP_ALIGN.CENTER)

# ======================== SLIDE 4: REBRANDING ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
slide_title(slide, "REBRANDING", "ASTraM Event Congestion Command to Grid Unlock Command")
add_rr(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.2), DARK_CARD)
tb(slide, Inches(1.1), Inches(2.0), Inches(5.0), Inches(0.3), "Before: ASTraM", 14, SLATE)
tb(slide, Inches(1.1), Inches(2.4), Inches(5.0), Inches(0.6), "ASTraM Event Congestion Command", 20, WHITE, True)
mt(slide, Inches(1.1), Inches(3.2), Inches(5.0), Inches(2.5), ["ASTraM logo and title across all pages", "ASTraM footer and copyright notice", "astram-events CSV download prefix", "ASTraM branding in console messages"], 13, SLATE)
add_rr(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.2), DARK_CARD)
tb(slide, Inches(7.1), Inches(2.0), Inches(5.0), Inches(0.3), "After: Grid Unlock", 14, SKY_BLUE)
tb(slide, Inches(7.1), Inches(2.4), Inches(5.0), Inches(0.6), "Grid Unlock Command", 20, WHITE, True)
mt(slide, Inches(7.1), Inches(3.2), Inches(5.0), Inches(2.5), ["Grid Unlock logo (letter G in navy circle)", "Updated footer and copyright notice", "grid-unlock CSV download prefix", "Updated console messages and page title"], 13, SLATE)
tb(slide, Inches(5.7), Inches(3.5), Inches(1.5), Inches(0.4), ">", 32, SKY_BLUE, al=PP_ALIGN.CENTER)
add_r(slide, Inches(0.8), Inches(6.5), Inches(11.5), Pt(2), ROYAL_BLUE)
tb(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.3), "Consistent branding across 7 UI touchpoints: title, navbar, footer, copyright, CSV filename, console logs, and page metadata.", 11, SKY_BLUE)

# ======================== SLIDE 5: UI ENHANCEMENTS ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
top_bar(slide)
tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), "USER INTERFACE ENHANCEMENTS", 28, DARK_TEXT, True)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4), "Premium SaaS landing page with comprehensive dashboard improvements", 14, GRAY_TEXT)
for idx, (title, desc) in enumerate([("Corridor Selector","Two-level dropdown: category selector (Corridor/Non-corridor) followed by road selector. Backend API returns categorized data."),("Input Validation","Hour: 1-12 clamp. Attendance: max 10,000. Officers: max 500. Barricades: max 1,000. Patrols: max 100. Notes: 200-word limit with live counter."),("Premium Landing Page","Glassmorphism cards with hover lift. Scroll-triggered reveal animations. Live digital clock. Six feature cards. Gradient hero section."),("Feedback Log Filters","Date from/to, Hour from/to dropdowns, Day of week, Year (auto-populated). Apply button disabled until filter selected. Newest-first sort."),("Ground Truth Form","Actual Duration editable with type=number, range 0-12 step 0.5, and AM/PM selector. Default changed from 1 to 0 (neutral baseline)."),("Filter Visibility Fix","Filter select elements use dark background (#1E293B) with white text. Hour from/to broken template literals replaced with JavaScript function.")]):
    col, row = idx % 3, idx // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.8 + row * 2.6)
    add_rr(slide, x, y, Inches(3.8), Inches(2.2), OFF_WHITE_BG)
    tb(slide, x + Inches(0.25), y + Inches(0.2), Inches(3.3), Inches(0.3), title, 13, DARK_TEXT, True)
    tb(slide, x + Inches(0.25), y + Inches(0.6), Inches(3.3), Inches(1.4), desc, 11, GRAY_TEXT)

# ======================== SLIDE 6: FEEDBACK LOG ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
slide_title(slide, "FEEDBACK LOG & DATA QUALITY", "Comprehensive feedback tracking with improved sort, filter, and display")
add_rr(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5), DARK_CARD)
tb(slide, Inches(1.1), Inches(2.0), Inches(5.3), Inches(0.3), "Filtering and Sorting", 16, ROYAL_BLUE, True)
mt(slide, Inches(1.1), Inches(2.5), Inches(5.3), Inches(3.5), ["Date range filter with from/to date pickers", "Hour from/to dropdowns (12-hour format)", "Day of week filter (Sunday through Saturday)", "Year filter auto-populated from log data", "Apply button disabled until filter selected", "Sort by Default, Timestamp, Type, Cause, Corridor", "Default sort order: newest first", "Chronological sort using Date.getTime()"], 13, SLATE)
add_rr(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5), DARK_CARD)
tb(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.3), "Display and Data Quality", 16, ROYAL_BLUE, True)
mt(slide, Inches(7.1), Inches(2.5), Inches(5.3), Inches(3.5), ["Paired columns: Predicted vs Actual for all metrics", "Severity, Duration, Officers, Barricades, Patrols", "tree_fall added to event cause dropdown", "Actual Duration: editable field with AM/PM selector", "Default duration set to 0 (neutral baseline)", "12 AM/PM display fix in timeline", "White text on dark background for all filter selects", "CSS variables for consistent theming"], 13, SLATE)
add_r(slide, Inches(0.8), Inches(6.7), Inches(11.5), Pt(2), ROYAL_BLUE)
tb(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.3), "All filters and sort options auto-apply on change via re-render triggers.", 11, SKY_BLUE)

# ======================== SLIDE 7: MODEL UPDATES ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
top_bar(slide)
tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), "MACHINE LEARNING MODEL UPDATES", 28, DARK_TEXT, True)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4), "All models retrained with 147-corridor dataset. No InconsistentVersionWarning.", 14, GRAY_TEXT)
for i, (nm, al, desc, met) in enumerate([("closure_model","Random Forest Classifier","Predicts if road closure is required","84.6% accuracy"),("duration_model","Random Forest Classifier","3-class: short/medium/extended","56.9% accuracy"),("severity_model","Random Forest Classifier","3-class: High/Medium/Low","76.3% accuracy"),("officers_model","Random Forest Regressor","Predicts number of officers","MAE 2.33"),("barricades_model","Random Forest Regressor","Predicts number of barricades","MAE 1.74"),("patrols_model","Random Forest Regressor","Predicts number of patrol vehicles","MAE 0.68")]):
    col, row = i % 3, i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.8 + row * 2.5)
    add_rr(slide, x, y, Inches(3.8), Inches(2.1), OFF_WHITE_BG)
    tb(slide, x + Inches(0.25), y + Inches(0.15), Inches(3.3), Inches(0.3), nm, 14, DARK_TEXT, True)
    tb(slide, x + Inches(0.25), y + Inches(0.5), Inches(1.8), Inches(0.4), al, 11, GRAY_TEXT)
    tb(slide, x + Inches(0.25), y + Inches(0.9), Inches(1.8), Inches(0.4), desc, 11, GRAY_TEXT)
    tb(slide, x + Inches(2.0), y + Inches(0.6), Inches(1.6), Inches(0.9), met, 14, ROYAL_BLUE, True, PP_ALIGN.CENTER)

# ======================== SLIDE 8: BACKEND PIPELINE ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
slide_title(slide, "BACKEND & DATA PIPELINE UPDATES", "Corresponding changes across all backend modules")
add_rr(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5), DARK_CARD)
tb(slide, Inches(1.1), Inches(2.0), Inches(5.3), Inches(0.3), "Module Updates", 16, ROYAL_BLUE, True)
mt(slide, Inches(1.1), Inches(2.5), Inches(5.3), Inches(3.5), ["src/noncorridor_roads.py: Created with NAMED_CORRIDORS set and 125 NON_CORRIDOR_ROADS entries", "src/spatial_engine.py: CORRIDOR_CENTROIDS updated to 147 entries for map markers", "src/corridor_network.py: 147 graph nodes with connectivity edges", "src/resource_optimizer.py: CORRIDOR_RADIUS_KM with 147 entries (0.8/2.5/5 km tiers)", "src/retrain_pipeline.py: Float cast fix for numpy type compatibility"], 13, SLATE)
add_rr(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5), DARK_CARD)
tb(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.3), "API and Data Changes", 16, ROYAL_BLUE, True)
mt(slide, Inches(7.1), Inches(2.5), Inches(5.3), Inches(3.5), ["/api/corridors returns categorized: named(21)+noncorridor(125)", "outputs/processed_theme3.csv regenerated with 147 corridors", "CSV re-read cache added to avoid stale data", "Leaflet CSS link added, map height set to 300px", "Field path fixes: road_closure, resources, routes, junctions"], 13, SLATE)
add_r(slide, Inches(0.8), Inches(6.7), Inches(11.5), Pt(2), ROYAL_BLUE)
tb(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.3), "6 backend modules updated to support 147-corridor architecture.", 11, SKY_BLUE)

# ======================== SLIDE 9: BUG FIXES ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
top_bar(slide)
tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), "BUG FIXES & CODE CLEANUP", 28, DARK_TEXT, True)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4), "Reliability and maintainability improvements across the codebase", 14, GRAY_TEXT)
add_rr(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.2), OFF_WHITE_BG)
tb(slide, Inches(1.1), Inches(2.0), Inches(5.3), Inches(0.3), "Stability Fixes", 15, DARK_TEXT, True)
mt(slide, Inches(1.1), Inches(2.5), Inches(5.3), Inches(3.2), ["Race conditions in retraining pipeline eliminated", "duration_bucket None/NaN handling fixed for edge cases", "retrain_pipeline.py: numpy float cast fix", "12 AM/PM display fix in timeline rendering", "CSV re-read cache to prevent stale data serving", "Orphaned temp files cleaned automatically"], 13, GRAY_TEXT)
add_rr(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.2), OFF_WHITE_BG)
tb(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.3), "Code Cleanup", 15, DARK_TEXT, True)
mt(slide, Inches(7.1), Inches(2.5), Inches(5.3), Inches(3.2), ["Dead code removed from feedback_logger.py (12 unreachable lines)", "Stale column ref timestamp to log_timestamp in delete", "priority_model.joblib removed (rule-based, not ML)", "Screenshot*.png added to .gitignore", "Duplicate Dashboard navigation button removed", "Testimonials section removed for cleaner UX"], 13, GRAY_TEXT)

# ======================== SLIDE 10: KEY METRICS ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
slide_title(slide, "KEY METRICS", "Quantifiable improvements across the platform")
for i, (v, l, d) in enumerate([("147","Total Corridors","21 named + 125 extracted + 1 fallback"),("125","Non-Corridor Roads","Extracted from 562 unique names"),("7,532","Processed Rows","Full dataset with 147 corridors"),("6","ML Models","All retrained with new corridor data"),("84.6%","Closure Accuracy","Best classification model"),("76.3%","Severity Accuracy","3-class severity prediction")]):
    col, row = i % 3, i // 3
    x = Inches(0.8 + col * 4.1); y = Inches(1.8 + row * 2.5)
    add_rr(slide, x, y, Inches(3.8), Inches(2.1), DARK_CARD)
    tb(slide, x + Inches(0.25), y + Inches(0.2), Inches(3.3), Inches(0.5), v, 26, ROYAL_BLUE, True)
    tb(slide, x + Inches(0.25), y + Inches(0.75), Inches(3.3), Inches(0.3), l, 13, WHITE, True)
    tb(slide, x + Inches(0.25), y + Inches(1.1), Inches(3.3), Inches(0.8), d, 12, SLATE)

# ======================== SLIDE 11: THANK YOU ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_r(slide, Inches(0), Inches(2.5), Inches(13.333), Pt(5), ROYAL_BLUE)
tb(slide, Inches(1), Inches(1.7), Inches(11), Inches(0.9), "Thank You", 44, WHITE, True, PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.5), "Grid Unlock Command  |  Event-Driven Congestion (Planned & Unplanned)", 15, SKY_BLUE, al=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.4), "Built with Flask, scikit-learn, Leaflet.js  |  6 ML Models  |  147 Bangalore Corridors", 13, SLATE, al=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.4), "Team: The ML Mavericks", 17, WHITE, al=PP_ALIGN.CENTER)
tb(slide, Inches(2), Inches(5.1), Inches(5), Inches(0.3), "Eaknath M S (Team Leader)", 13, SLATE, al=PP_ALIGN.CENTER)
tb(slide, Inches(7), Inches(5.1), Inches(5), Inches(0.3), "Ashish D", 13, SLATE, al=PP_ALIGN.CENTER)
tb(slide, Inches(2), Inches(5.5), Inches(5), Inches(0.3), "Sabarishh P", 13, SLATE, al=PP_ALIGN.CENTER)
tb(slide, Inches(7), Inches(5.5), Inches(5), Inches(0.3), "Prasanth Kanaga Sabai S", 13, SLATE, al=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.3), "Predict  |  Deploy  |  Learn", 13, ROYAL_BLUE, al=PP_ALIGN.CENTER)

prs.save("Grid_Unlock_Updates.pptx")
print("Saved Grid_Unlock_Updates.pptx - 11 slides")
